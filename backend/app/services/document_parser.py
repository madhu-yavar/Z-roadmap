import csv
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation


def _safe_read_text(path: Path) -> str:
    try:
        return path.read_text(encoding="utf-8", errors="ignore")
    except Exception:
        return ""


def _clean_cell_value(cell: object) -> str:
    if cell is None:
        return ""
    return " ".join(str(cell).strip().split())


def _is_header_like(values: list[str]) -> bool:
    if not values:
        return False
    text = " ".join(values).lower()
    hints = ("activity", "task", "description", "workstream", "owner", "role", "lane", "module")
    return any(h in text for h in hints)


def extract_document_units(file_path: str, file_type: str) -> list[dict]:
    path = Path(file_path)
    ext = file_type.lower().strip(".") or path.suffix.lower().strip(".")

    units: list[dict] = []

    if ext == "pdf":
        reader = PdfReader(str(path))
        for idx, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if not text:
                continue
            lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
            if not lines:
                units.append({"ref": f"page:{idx}", "text": text})
                continue
            for line_idx, line in enumerate(lines, start=1):
                units.append({"ref": f"page:{idx}:line:{line_idx}", "text": line})
        return units

    if ext in {"docx", "doc"}:
        doc = DocxDocument(str(path))
        seen_texts: set[str] = set()
        for idx, p in enumerate(doc.paragraphs, start=1):
            text = (p.text or "").strip()
            if text and text not in seen_texts:
                units.append({"ref": f"paragraph:{idx}", "text": text})
                seen_texts.add(text)

        # Many BRD/RFP Word files keep core requirements inside tables.
        for table_idx, table in enumerate(doc.tables, start=1):
            for row_idx, row in enumerate(table.rows, start=1):
                row_values = []
                for cell in row.cells:
                    cell_text = (cell.text or "").strip()
                    if cell_text:
                        row_values.append(" ".join(cell_text.split()))
                if not row_values:
                    continue
                row_text = " | ".join(row_values)
                if row_text in seen_texts:
                    continue
                units.append({"ref": f"table:{table_idx}:row:{row_idx}", "text": row_text})
                seen_texts.add(row_text)

        # Include header/footer text where objectives or scope notes may appear.
        for section_idx, section in enumerate(doc.sections, start=1):
            for para_idx, p in enumerate(section.header.paragraphs, start=1):
                text = (p.text or "").strip()
                if text and text not in seen_texts:
                    units.append({"ref": f"header:{section_idx}:paragraph:{para_idx}", "text": text})
                    seen_texts.add(text)
            for para_idx, p in enumerate(section.footer.paragraphs, start=1):
                text = (p.text or "").strip()
                if text and text not in seen_texts:
                    units.append({"ref": f"footer:{section_idx}:paragraph:{para_idx}", "text": text})
                    seen_texts.add(text)
        return units

    if ext in {"ppt", "pptx"}:
        ppt = Presentation(str(path))
        for slide_idx, slide in enumerate(ppt.slides, start=1):
            for shape_idx, shape in enumerate(slide.shapes, start=1):
                text = getattr(shape, "text", "")
                text = (text or "").strip()
                if text:
                    units.append({"ref": f"slide:{slide_idx}:shape:{shape_idx}", "text": text})
        return units

    if ext in {"xlsx", "xls"}:
        try:
            wb = load_workbook(str(path), data_only=True)
        except Exception:
            # Keep graceful fallback for legacy/unsupported spreadsheet variants.
            raw = _safe_read_text(path)
            for idx, line in enumerate(raw.splitlines(), start=1):
                line = line.strip()
                if line:
                    units.append({"ref": f"line:{idx}", "text": line})
            return units

        for ws in wb.worksheets[:6]:
            headers: list[str] = []
            for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=240, max_col=20, values_only=True), start=1):
                row_values = [_clean_cell_value(cell) for cell in row]
                trimmed = [v for v in row_values if v]
                if not trimmed:
                    continue
                if row_idx == 1 and _is_header_like(trimmed):
                    headers = trimmed
                    units.append({"ref": f"sheet:{ws.title}:header", "text": " | ".join(headers)})
                    continue

                if headers:
                    mapped = []
                    for col_idx, val in enumerate(row_values):
                        if not val:
                            continue
                        key = headers[col_idx] if col_idx < len(headers) and headers[col_idx] else f"col_{col_idx+1}"
                        mapped.append(f"{key}: {val}")
                    if mapped:
                        units.append({"ref": f"sheet:{ws.title}:row:{row_idx}", "text": " | ".join(mapped)})
                        continue
                units.append({"ref": f"sheet:{ws.title}:row:{row_idx}", "text": " | ".join(trimmed)})
        return units

    if ext == "csv":
        raw = _safe_read_text(path)
        if not raw:
            return units
        sample = raw[:2048]
        try:
            dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
        except Exception:
            dialect = csv.excel
        reader = csv.reader(raw.splitlines(), dialect=dialect)
        rows = list(reader)
        headers: list[str] = []
        if rows:
            first = [_clean_cell_value(x) for x in rows[0]]
            if _is_header_like([x for x in first if x]):
                headers = first
                units.append({"ref": "csv:header", "text": " | ".join([h for h in headers if h])})
                start_idx = 2
            else:
                start_idx = 1
            for row_idx, row in enumerate(rows[start_idx - 1 : min(len(rows), 260)], start=start_idx):
                row_values = [_clean_cell_value(x) for x in row]
                trimmed = [v for v in row_values if v]
                if not trimmed:
                    continue
                if headers:
                    mapped = []
                    for col_idx, val in enumerate(row_values):
                        if not val:
                            continue
                        key = headers[col_idx] if col_idx < len(headers) and headers[col_idx] else f"col_{col_idx+1}"
                        mapped.append(f"{key}: {val}")
                    if mapped:
                        units.append({"ref": f"csv:row:{row_idx}", "text": " | ".join(mapped)})
                        continue
                units.append({"ref": f"csv:row:{row_idx}", "text": " | ".join(trimmed)})
        return units

    raw = _safe_read_text(path)
    for idx, line in enumerate(raw.splitlines(), start=1):
        line = line.strip()
        if line:
            units.append({"ref": f"line:{idx}", "text": line})
    return units


def extract_document_text(file_path: str, file_type: str) -> str:
    units = extract_document_units(file_path=file_path, file_type=file_type)
    return "\n".join(unit["text"] for unit in units)
